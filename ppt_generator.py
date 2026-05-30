"""
PPT Generator - PowerPoint 演示文稿自动生成模块
功能：根据大纲自动生成幻灯片，支持图表插入和样式美化
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
import pandas as pd
from typing import List, Dict, Optional


class PPTGenerator:
    """PPT 生成器类"""
    
    def __init__(self, template_path: Optional[str] = None):
        """
        初始化 PPT 生成器
        
        Args:
            template_path: 模板文件路径，None 则使用默认模板
        """
        if template_path:
            self.prs = Presentation(template_path)
        else:
            self.prs = Presentation()
        
        # 默认配色方案
        self.color_scheme = {
            'primary': RGBColor(41, 128, 185),    # 蓝色
            'secondary': RGBColor(142, 68, 173),  # 紫色
            'accent': RGBColor(230, 126, 34),     # 橙色
            'text': RGBColor(44, 62, 80),         # 深灰
            'background': RGBColor(255, 255, 255) # 白色
        }
    
    def set_color_scheme(self, scheme_name: str):
        """
        设置配色方案
        
        Args:
            scheme_name: 配色方案名称 ('professional', 'vibrant', 'minimal', 'dark')
        """
        schemes = {
            'professional': {
                'primary': RGBColor(41, 128, 185),
                'secondary': RGBColor(52, 73, 94),
                'accent': RGBColor(230, 126, 34),
                'text': RGBColor(44, 62, 80),
                'background': RGBColor(255, 255, 255)
            },
            'vibrant': {
                'primary': RGBColor(231, 76, 60),
                'secondary': RGBColor(52, 152, 219),
                'accent': RGBColor(241, 196, 15),
                'text': RGBColor(44, 62, 80),
                'background': RGBColor(255, 255, 255)
            },
            'minimal': {
                'primary': RGBColor(0, 0, 0),
                'secondary': RGBColor(128, 128, 128),
                'accent': RGBColor(0, 128, 128),
                'text': RGBColor(0, 0, 0),
                'background': RGBColor(255, 255, 255)
            },
            'dark': {
                'primary': RGBColor(52, 152, 219),
                'secondary': RGBColor(155, 89, 182),
                'accent': RGBColor(241, 196, 15),
                'text': RGBColor(236, 240, 241),
                'background': RGBColor(44, 62, 80)
            }
        }
        
        if scheme_name in schemes:
            self.color_scheme = schemes[scheme_name]
        else:
            raise ValueError(f"未知的配色方案：{scheme_name}")
    
    def add_title_slide(self, title: str, subtitle: str = "") -> None:
        """
        添加标题幻灯片
        
        Args:
            title: 标题文本
            subtitle: 副标题文本
        """
        slide_layout = self.prs.slide_layouts[0]  # Title Slide
        slide = self.prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]
        
        title_shape.text = title
        subtitle_shape.text = subtitle
        
        # 设置标题样式
        for paragraph in title_shape.text_frame.paragraphs:
            paragraph.font.size = Pt(44)
            paragraph.font.bold = True
            paragraph.font.color.rgb = self.color_scheme['primary']
        
        # 设置副标题样式
        for paragraph in subtitle_shape.text_frame.paragraphs:
            paragraph.font.size = Pt(24)
            paragraph.font.color.rgb = self.color_scheme['secondary']
    
    def add_content_slide(self, title: str, content: List[str], 
                         bullet_level: int = 0) -> None:
        """
        添加内容幻灯片
        
        Args:
            title: 幻灯片标题
            content: 内容列表（项目符号）
            bullet_level: 项目符号层级
        """
        slide_layout = self.prs.slide_layouts[1]  # Title and Content
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 设置标题
        title_shape = slide.shapes.title
        title_shape.text = title
        
        for paragraph in title_shape.text_frame.paragraphs:
            paragraph.font.size = Pt(36)
            paragraph.font.bold = True
            paragraph.font.color.rgb = self.color_scheme['primary']
        
        # 设置内容
        content_shape = slide.placeholders[1]
        text_frame = content_shape.text_frame
        text_frame.clear()
        
        for i, item in enumerate(content):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            p.text = item
            p.level = bullet_level
            p.font.size = Pt(18)
            p.font.color.rgb = self.color_scheme['text']
            p.space_after = Pt(10)
    
    def add_chart_slide(self, title: str, chart_type: str, 
                       data: pd.DataFrame, x_column: str, 
                       y_columns: List[str]) -> None:
        """
        添加图表幻灯片
        
        Args:
            title: 幻灯片标题
            chart_type: 图表类型 ('bar', 'line', 'pie', 'area')
            data: pandas DataFrame 数据
            x_column: X 轴列名
            y_columns: Y 轴列名列表
        """
        slide_layout = self.prs.slide_layouts[5]  # Blank
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 添加标题
        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(9)
        height = Inches(1)
        
        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.text = title
        
        for paragraph in title_frame.paragraphs:
            paragraph.font.size = Pt(32)
            paragraph.font.bold = True
            paragraph.font.color.rgb = self.color_scheme['primary']
        
        # 准备图表数据
        chart_data = CategoryChartData()
        chart_data.categories = data[x_column].tolist()
        
        for col in y_columns:
            chart_data.add_series(col, data[col].tolist())
        
        # 添加图表
        left = Inches(1)
        top = Inches(1.5)
        width = Inches(8)
        height = Inches(5)
        
        chart_types = {
            'bar': XL_CHART_TYPE.COLUMN_CLUSTERED,
            'line': XL_CHART_TYPE.LINE,
            'pie': XL_CHART_TYPE.PIE,
            'area': XL_CHART_TYPE.AREA
        }
        
        if chart_type not in chart_types:
            raise ValueError(f"不支持的图表类型：{chart_type}")
        
        chart = slide.shapes.add_chart(
            chart_types[chart_type], left, top, width, height, chart_data
        ).chart
        
        # 设置图表样式
        chart.has_legend = True
        chart.legend.include_in_layout = False
        
        if chart_type != 'pie':
            chart.has_category_axis = True
            chart.has_value_axis = True
    
    def add_image_slide(self, title: str, image_path: str, 
                       caption: str = "") -> None:
        """
        添加图片幻灯片
        
        Args:
            title: 幻灯片标题
            image_path: 图片文件路径
            caption: 图片说明文字
        """
        slide_layout = self.prs.slide_layouts[5]  # Blank
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 添加标题
        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(9)
        height = Inches(1)
        
        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.text = title
        
        for paragraph in title_frame.paragraphs:
            paragraph.font.size = Pt(32)
            paragraph.font.bold = True
            paragraph.font.color.rgb = self.color_scheme['primary']
        
        # 添加图片
        left = Inches(1)
        top = Inches(1.5)
        width = Inches(8)
        height = Inches(5)
        
        slide.shapes.add_picture(image_path, left, top, width=width, height=height)
        
        # 添加说明文字
        if caption:
            caption_top = Inches(6.8)
            caption_box = slide.shapes.add_textbox(left, caption_top, width, Inches(0.5))
            caption_frame = caption_box.text_frame
            caption_frame.text = caption
            
            for paragraph in caption_frame.paragraphs:
                paragraph.font.size = Pt(14)
                paragraph.font.color.rgb = self.color_scheme['secondary']
                paragraph.alignment = PP_ALIGN.CENTER
    
    def add_section_header(self, section_title: str) -> None:
        """
        添加章节标题页
        
        Args:
            section_title: 章节标题
        """
        slide_layout = self.prs.slide_layouts[5]  # Blank
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 添加背景形状
        left = Inches(0)
        top = Inches(0)
        width = Inches(10)
        height = Inches(7.5)
        
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.color_scheme['primary']
        shape.line.fill.background()
        
        # 添加标题
        text_box = slide.shapes.add_textbox(
            Inches(1), Inches(2.5), Inches(8), Inches(2)
        )
        text_frame = text_box.text_frame
        text_frame.text = section_title
        
        for paragraph in text_frame.paragraphs:
            paragraph.font.size = Pt(48)
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(255, 255, 255)
            paragraph.alignment = PP_ALIGN.CENTER
    
    def generate_from_outline(self, outline: Dict) -> None:
        """
        根据大纲生成完整 PPT
        
        Args:
            outline: 大纲字典，格式如下：
            {
                'title': '演示文稿标题',
                'subtitle': '副标题',
                'sections': [
                    {
                        'title': '章节标题',
                        'slides': [
                            {
                                'type': 'content',
                                'title': '幻灯片标题',
                                'content': ['要点 1', '要点 2']
                            },
                            {
                                'type': 'chart',
                                'title': '图表标题',
                                'chart_type': 'bar',
                                'data': dataframe,
                                'x_column': 'x',
                                'y_columns': ['y1', 'y2']
                            }
                        ]
                    }
                ]
            }
        """
        # 添加标题页
        self.add_title_slide(outline['title'], outline.get('subtitle', ''))
        
        # 添加各章节
        for section in outline.get('sections', []):
            # 添加章节标题页
            if 'title' in section:
                self.add_section_header(section['title'])
            
            # 添加幻灯片
            for slide_data in section.get('slides', []):
                slide_type = slide_data.get('type', 'content')
                
                if slide_type == 'content':
                    self.add_content_slide(
                        slide_data['title'],
                        slide_data.get('content', [])
                    )
                elif slide_type == 'chart':
                    self.add_chart_slide(
                        slide_data['title'],
                        slide_data['chart_type'],
                        slide_data['data'],
                        slide_data['x_column'],
                        slide_data['y_columns']
                    )
                elif slide_type == 'image':
                    self.add_image_slide(
                        slide_data['title'],
                        slide_data['image_path'],
                        slide_data.get('caption', '')
                    )
    
    def save(self, output_path: str) -> None:
        """
        保存 PPT 文件
        
        Args:
            output_path: 输出文件路径
        """
        self.prs.save(output_path)
    
    def save_as_pdf(self, output_path: str) -> None:
        """
        保存为 PDF 文件（需要安装 comtypes）
        
        Args:
            output_path: 输出 PDF 文件路径
        """
        import os
        import comtypes.client
        
        # 先保存为 PPTX
        temp_pptx = output_path.replace('.pdf', '_temp.pptx')
        self.save(temp_pptx)
        
        # 转换为 PDF
        powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
        powerpoint.Visible = 1
        
        presentation = powerpoint.Presentations.Open(
            os.path.abspath(temp_pptx),
            WithWindow=False
        )
        
        presentation.SaveAs(
            os.path.abspath(output_path),
            32  # ppSaveAsPDF
        )
        
        presentation.Close()
        powerpoint.Quit()
        
        # 删除临时文件
        os.remove(temp_pptx)
    
    def get_slide_count(self) -> int:
        """获取幻灯片数量"""
        return len(self.prs.slides)


def create_sample_outline() -> Dict:
    """创建示例大纲"""
    import pandas as pd
    
    # 创建示例数据
    df = pd.DataFrame({
        '季度': ['Q1', 'Q2', 'Q3', 'Q4'],
        '销售额': [150, 200, 180, 250],
        '利润': [30, 45, 40, 60]
    })
    
    outline = {
        'title': '2024 年度业务报告',
        'subtitle': '业绩总结与展望',
        'sections': [
            {
                'title': '第一部分：业绩概览',
                'slides': [
                    {
                        'type': 'content',
                        'title': '年度亮点',
                        'content': [
                            '销售额同比增长 25%',
                            '新客户增长 40%',
                            '市场份额提升至 15%',
                            '客户满意度达 95%'
                        ]
                    },
                    {
                        'type': 'chart',
                        'title': '季度销售趋势',
                        'chart_type': 'bar',
                        'data': df,
                        'x_column': '季度',
                        'y_columns': ['销售额', '利润']
                    }
                ]
            },
            {
                'title': '第二部分：未来规划',
                'slides': [
                    {
                        'type': 'content',
                        'title': '2025 年目标',
                        'content': [
                            '销售额目标：增长 30%',
                            '拓展 5 个新市场',
                            '推出 3 款新产品',
                            '团队规模扩大 50%'
                        ]
                    }
                ]
            }
        ]
    }
    
    return outline


if __name__ == "__main__":
    # 测试代码
    print("PPT 生成器测试开始")

    generator = PPTGenerator()
    generator.set_color_scheme('professional')
    
    # 使用示例大纲生成 PPT
    outline = create_sample_outline()
    generator.generate_from_outline(outline)
    
    # 保存文件
    output_path = "test_presentation.pptx"
    generator.save(output_path)
    
    print(f"✓ PPT 生成成功：{output_path}")
    print(f"✓ 幻灯片数量：{generator.get_slide_count()}")
    print("PPT 生成器测试完成")
