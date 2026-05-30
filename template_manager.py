"""
Template Manager - PPT 模板管理模块
功能：管理多种 PPT 模板，支持模板预览和应用
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from typing import Dict, List, Optional
import json


class TemplateManager:
    """PPT 模板管理器"""
    
    def __init__(self, template_dir: str = "templates"):
        """
        初始化模板管理器
        
        Args:
            template_dir: 模板目录路径
        """
        self.template_dir = template_dir
        self.templates = {}
        self.load_templates()
    
    def load_templates(self):
        """加载所有可用模板"""
        if not os.path.exists(self.template_dir):
            os.makedirs(self.template_dir)
            # 创建默认模板
            self.create_default_templates()
        
        # 扫描模板文件
        for file in os.listdir(self.template_dir):
            if file.endswith('.pptx'):
                template_name = file.replace('.pptx', '')
                template_path = os.path.join(self.template_dir, file)
                self.templates[template_name] = template_path
        
        # 如果没有模板，创建默认模板
        if not self.templates:
            self.create_default_templates()
            self.load_templates()
    
    def create_default_templates(self):
        """创建默认模板"""
        # 专业蓝色模板
        self._create_professional_template()
        
        # 活力红色模板
        self._create_vibrant_template()
        
        # 极简黑白模板
        self._create_minimal_template()
        
        # 深色模式模板
        self._create_dark_template()
    
    def _create_professional_template(self):
        """创建专业蓝色模板"""
        prs = Presentation()
        
        # 设置幻灯片母版
        master = prs.slide_master
        background = master.background
        background.fill.solid()
        background.fill.fore_color.rgb = RGBColor(255, 255, 255)
        
        # 保存模板
        template_path = os.path.join(self.template_dir, 'professional.pptx')
        prs.save(template_path)
        
        # 保存模板元数据
        self._save_template_metadata('professional', {
            'name': '专业蓝色',
            'description': '适合商务报告、正式演示',
            'primary_color': '#2980B9',
            'secondary_color': '#34495E',
            'accent_color': '#E67E22',
            'style': 'professional'
        })
    
    def _create_vibrant_template(self):
        """创建活力红色模板"""
        prs = Presentation()
        
        # 设置幻灯片母版
        master = prs.slide_master
        background = master.background
        background.fill.solid()
        background.fill.fore_color.rgb = RGBColor(255, 255, 255)
        
        # 保存模板
        template_path = os.path.join(self.template_dir, 'vibrant.pptx')
        prs.save(template_path)
        
        # 保存模板元数据
        self._save_template_metadata('vibrant', {
            'name': '活力红色',
            'description': '适合产品发布、创意演示',
            'primary_color': '#E74C3C',
            'secondary_color': '#3498DB',
            'accent_color': '#F1C40F',
            'style': 'vibrant'
        })
    
    def _create_minimal_template(self):
        """创建极简黑白模板"""
        prs = Presentation()
        
        # 设置幻灯片母版
        master = prs.slide_master
        background = master.background
        background.fill.solid()
        background.fill.fore_color.rgb = RGBColor(255, 255, 255)
        
        # 保存模板
        template_path = os.path.join(self.template_dir, 'minimal.pptx')
        prs.save(template_path)
        
        # 保存模板元数据
        self._save_template_metadata('minimal', {
            'name': '极简黑白',
            'description': '适合学术报告、简洁演示',
            'primary_color': '#000000',
            'secondary_color': '#808080',
            'accent_color': '#008080',
            'style': 'minimal'
        })
    
    def _create_dark_template(self):
        """创建深色模式模板"""
        prs = Presentation()
        
        # 设置幻灯片母版
        master = prs.slide_master
        background = master.background
        background.fill.solid()
        background.fill.fore_color.rgb = RGBColor(44, 70, 80)
        
        # 保存模板
        template_path = os.path.join(self.template_dir, 'dark.pptx')
        prs.save(template_path)
        
        # 保存模板元数据
        self._save_template_metadata('dark', {
            'name': '深色模式',
            'description': '适合科技演示、晚间展示',
            'primary_color': '#3498DB',
            'secondary_color': '#9B59B6',
            'accent_color': '#F1C40F',
            'style': 'dark'
        })
    
    def _save_template_metadata(self, name: str, metadata: Dict):
        """保存模板元数据"""
        metadata_path = os.path.join(self.template_dir, f'{name}.json')
        with open(metadata_path, 'w', encoding='utf-8') as f:
            json.dump(metadata, f, ensure_ascii=False, indent=2)
    
    def get_template_list(self) -> List[Dict]:
        """
        获取所有模板列表
        
        Returns:
            模板信息列表
        """
        template_list = []
        
        for name, path in self.templates.items():
            metadata_path = os.path.join(self.template_dir, f'{name}.json')
            
            if os.path.exists(metadata_path):
                with open(metadata_path, 'r', encoding='utf-8') as f:
                    metadata = json.load(f)
            else:
                metadata = {
                    'name': name,
                    'description': '自定义模板',
                    'style': 'custom'
                }
            
            metadata['id'] = name
            metadata['path'] = path
            template_list.append(metadata)
        
        return template_list
    
    def get_template(self, template_name: str) -> Optional[str]:
        """
        获取模板文件路径
        
        Args:
            template_name: 模板名称
            
        Returns:
            模板文件路径，不存在则返回 None
        """
        return self.templates.get(template_name)
    
    def apply_template(self, template_name: str) -> Optional[Presentation]:
        """
        应用模板
        
        Args:
            template_name: 模板名称
            
        Returns:
            Presentation 对象，模板不存在则返回 None
        """
        template_path = self.get_template(template_name)
        
        if template_path and os.path.exists(template_path):
            return Presentation(template_path)
        
        return None
    
    def delete_template(self, template_name: str) -> bool:
        """
        删除模板
        
        Args:
            template_name: 模板名称
            
        Returns:
            是否删除成功
        """
        if template_name in self.templates:
            template_path = self.templates[template_name]
            metadata_path = os.path.join(self.template_dir, f'{template_name}.json')
            
            # 删除模板文件
            if os.path.exists(template_path):
                os.remove(template_path)
            
            # 删除元数据文件
            if os.path.exists(metadata_path):
                os.remove(metadata_path)
            
            # 从字典中移除
            del self.templates[template_name]
            
            return True
        
        return False
    
    def preview_template(self, template_name: str) -> Dict:
        """
        预览模板信息
        
        Args:
            template_name: 模板名称
            
        Returns:
            模板预览信息
        """
        template_path = self.get_template(template_name)
        
        if not template_path or not os.path.exists(template_path):
            return {'error': '模板不存在'}
        
        # 加载模板
        prs = Presentation(template_path)
        
        # 获取基本信息
        preview_info = {
            'name': template_name,
            'slide_count': len(prs.slides),
            'slide_layouts': len(prs.slide_layouts)
        }
        
        # 获取母版信息
        if prs.slide_masters:
            master = prs.slide_masters[0]
            preview_info['master_name'] = master.name
        
        return preview_info
    
    def create_custom_template(self, name: str, color_scheme: Dict, 
                              description: str = "") -> str:
        """
        创建自定义模板
        
        Args:
            name: 模板名称
            color_scheme: 配色方案字典
            description: 模板描述
            
        Returns:
            模板文件路径
        """
        prs = Presentation()
        
        # 设置幻灯片母版
        master = prs.slide_master
        background = master.background
        background.fill.solid()
        
        # 设置背景色
        if 'background' in color_scheme:
            bg_color = color_scheme['background']
            if isinstance(bg_color, str):
                bg_color = bg_color.lstrip('#')
                r, g, b = int(bg_color[0:2], 16), int(bg_color[2:4], 16), int(bg_color[4:6], 16)
                background.fill.fore_color.rgb = RGBColor(r, g, b)
            else:
                background.fill.fore_color.rgb = bg_color
        
        # 保存模板
        template_path = os.path.join(self.template_dir, f'{name}.pptx')
        prs.save(template_path)
        
        # 保存元数据
        metadata = {
            'name': name,
            'description': description or '自定义模板',
            'primary_color': color_scheme.get('primary', '#000000'),
            'secondary_color': color_scheme.get('secondary', '#808080'),
            'accent_color': color_scheme.get('accent', '#FF0000'),
            'style': 'custom'
        }
        
        self._save_template_metadata(name, metadata)
        
        # 更新模板列表
        self.templates[name] = template_path
        
        return template_path


def get_available_templates() -> List[Dict]:
    """获取所有可用模板"""
    manager = TemplateManager()
    return manager.get_template_list()


if __name__ == "__main__":
    # 测试代码
    print("模板管理器测试开始")

    manager = TemplateManager()
    
    # 获取模板列表
    templates = manager.get_template_list()
    print(f"✓ 可用模板数量：{len(templates)}")
    
    for template in templates:
        print(f"  - {template['name']}: {template['description']}")
    
    # 预览模板
    if templates:
        preview = manager.preview_template(templates[0]['id'])
        print(f"\n✓ 模板预览：{preview}")

    print("\n模板管理器测试完成")
